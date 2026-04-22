import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import PyPDF2

def extract_pdf_text(pdf_path):
    """Extract text from PDF file"""
    text_content = []
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text.strip():
                    text_content.append(text)
        return text_content
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return []

def clone_presentation_style(template_path):
    """Clone the template presentation to get its style and layout"""
    try:
        template_prs = Presentation(template_path)
        return template_prs
    except Exception as e:
        print(f"Error loading template: {e}")
        return None

def create_presentation_from_pdf(pdf_path, template_path, output_path):
    """Create a new presentation based on PDF content and template style"""
    
    # Extract PDF content
    pdf_text = extract_pdf_text(pdf_path)
    if not pdf_text:
        print("No text found in PDF")
        return False
    
    # Clone template
    prs = clone_presentation_style(template_path)
    if not prs:
        prs = Presentation()
    
    # Get template layout if available
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    
    # Add title slide
    slide = prs.slides.add_slide(title_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    
    title.text = "Analysis Report"
    if subtitle:
        subtitle.text = "Generated from Final Report PDF"
    
    # Process PDF content and add to slides
    current_text = ""
    slide_count = 0
    
    for page_text in pdf_text:
        # Split text into chunks for slides
        paragraphs = page_text.split('\n')
        
        for para in paragraphs:
            if len(para.strip()) > 0:
                current_text += para.strip() + "\n"
                
                # Create a new slide when content reaches certain length
                if len(current_text) > 500 or slide_count > 20:
                    slide = prs.slides.add_slide(content_layout)
                    
                    # Add title
                    if slide.shapes.title:
                        slide.shapes.title.text = f"Content Slide {slide_count + 1}"
                    
                    # Add content
                    if len(slide.placeholders) > 1:
                        body_shape = slide.placeholders[1]
                        text_frame = body_shape.text_frame
                        text_frame.clear()
                        
                        # Add text with formatting
                        for line in current_text.split('\n')[:10]:  # Limit lines per slide
                            if line.strip():
                                p = text_frame.add_paragraph()
                                p.text = line.strip()[:100]  # Limit line length
                                p.level = 0
                                p.font.size = Pt(11)
                    
                    current_text = ""
                    slide_count += 1
    
    # Add remaining content
    if current_text.strip():
        slide = prs.slides.add_slide(content_layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"Final Content"
        
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()
            
            for line in current_text.split('\n')[:10]:
                if line.strip():
                    p = text_frame.add_paragraph()
                    p.text = line.strip()[:100]
                    p.level = 0
                    p.font.size = Pt(11)
    
    # Save presentation
    try:
        prs.save(output_path)
        print(f"✅ Presentation created successfully: {output_path}")
        return True
    except Exception as e:
        print(f"❌ Error saving presentation: {e}")
        return False

def main():
    # File paths
    pdf_path = "Final_report.pdf"
    template_path = "_DisasterDetection_PPT 23bcs0160.pptx"
    output_path = "Analysis_Report.pptx"
    
    # Check if files exist
    if not os.path.exists(pdf_path):
        print(f"❌ PDF file not found: {pdf_path}")
        return
    
    if not os.path.exists(template_path):
        print(f"⚠️  Template file not found: {template_path}")
        print("Creating presentation without template style")
    
    print("🚀 Starting PDF to PPTX conversion...")
    print(f"📄 Reading PDF: {pdf_path}")
    print(f"🎨 Using template: {template_path}")
    print(f"💾 Output will be saved to: {output_path}")
    
    success = create_presentation_from_pdf(pdf_path, template_path, output_path)
    
    if success:
        print(f"\n✅ Done! Download your presentation here:")
        print(f"📥 File: {output_path}")
        print(f"🔗 GitHub URL: https://github.com/jiohubconnect-ai/review/blob/main/{output_path}")

if __name__ == "__main__":
    main()